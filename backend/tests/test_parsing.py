import io
import json
import shutil
import subprocess
import uuid
import zipfile
from pathlib import Path

import httpx
import pytest
from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation
from pydantic import ValidationError

from robust_rag.parsing.base import FileMetadata, ParseError
from robust_rag.parsing.canonicalizer import Canonicalizer
from robust_rag.parsing.mineru import MinerUParser
from robust_rag.parsing.native import (
    ExcelParser,
    HtmlParser,
    LegacyOfficeParser,
    MarkdownParser,
    PlainTextParser,
    PowerPointParser,
    WordParser,
)
from robust_rag.parsing.router import ParserRouter
from robust_rag.parsing.schemas import BlockType, ParseArtifact, ParsedBlock


def metadata(path: Path, mime_type: str) -> FileMetadata:
    return FileMetadata(
        filename=path.name,
        mime_type=mime_type,
        file_size=path.stat().st_size,
        sha256="a" * 64,
    )


def test_text_markdown_and_html_are_structured_and_traceable(tmp_path: Path) -> None:
    text_path = tmp_path / "mixed.txt"
    text_path.write_text("第一段 Chinese.\n\nSecond paragraph.", encoding="utf-8")
    text = PlainTextParser().parse(text_path, metadata(text_path, "text/plain"))
    assert [block.original_text for block in text.blocks] == [
        "第一段 Chinese.",
        "Second paragraph.",
    ]
    assert text.blocks[1].source_locators[0].line_start == 3

    markdown_path = tmp_path / "guide.md"
    markdown_path.write_text(
        "# 标题 Title\n\n正文 text.\n\n- First\n- 第二项\n\n```python\nprint(1)\n```\n",
        encoding="utf-8",
    )
    markdown = MarkdownParser().parse(markdown_path, metadata(markdown_path, "text/markdown"))
    assert [block.block_type for block in markdown.blocks] == [
        BlockType.HEADING,
        BlockType.PARAGRAPH,
        BlockType.LIST,
        BlockType.LIST_ITEM,
        BlockType.LIST_ITEM,
        BlockType.CODE,
    ]
    assert markdown.blocks[3].parent_ref == markdown.blocks[2].ref
    assert markdown.blocks[0].source_locators[0].line_start == 1

    html_path = tmp_path / "page.html"
    html_path.write_text(
        "<html><head><title>企业手册</title></head><body><nav>丢弃</nav>"
        "<h1>Policy</h1><p>See <a href='/rule'>规则</a>.</p>"
        "<ul><li>First</li><li>第二项</li></ul>"
        "<table><tr><td>A</td><td>B</td></tr></table><footer>丢弃</footer></body></html>",
        encoding="utf-8",
    )
    html = HtmlParser().parse(html_path, metadata(html_path, "text/html"))
    assert html.title == "企业手册"
    assert "丢弃" not in " ".join(block.original_text for block in html.blocks)
    assert html.blocks[1].attributes["links"] == [{"text": "规则", "url": "/rule"}]
    assert html.blocks[1].source_locators[0].dom_path
    assert html.blocks[2].block_type is BlockType.LIST
    assert html.blocks[3].parent_ref == html.blocks[2].ref


def test_word_powerpoint_and_excel_native_parsers(tmp_path: Path) -> None:
    word_path = tmp_path / "handbook.docx"
    document = WordDocument()
    document.core_properties.title = "企业手册"
    document.add_heading("总则", level=1)
    document.add_paragraph("中英混合 policy")
    document.add_paragraph("第一项", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "名称"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "SLA"
    table.cell(1, 1).text = "99%"
    document.add_paragraph("表后说明")
    document.save(str(word_path))
    with zipfile.ZipFile(word_path, "a") as archive:
        archive.writestr(
            "word/footnotes.xml",
            "<w:footnotes xmlns:w='http://schemas.openxmlformats.org/wordprocessingml/2006/main'>"
            "<w:footnote w:id='1'><w:p><w:r><w:t>脚注 Footnote</w:t></w:r></w:p></w:footnote>"
            "</w:footnotes>",
        )
    word = WordParser().parse(word_path, metadata(word_path, "application/docx"))
    assert word.title == "企业手册"
    assert {block.block_type for block in word.blocks} >= {
        BlockType.HEADING,
        BlockType.PARAGRAPH,
        BlockType.LIST_ITEM,
        BlockType.TABLE,
        BlockType.FOOTNOTE,
    }
    table_position = next(
        index for index, block in enumerate(word.blocks) if block.block_type is BlockType.TABLE
    )
    assert word.blocks[0].source_locators[0].paragraph_index == 1
    assert word.blocks[table_position].source_locators[0].table_index == 1
    assert word.blocks[table_position + 1].original_text == "表后说明"
    assert word.blocks[-1].attributes["footnote_id"] == 1

    ppt_path = tmp_path / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "季度汇报"
    slide.placeholders[1].text = "Revenue 收入"
    rows, cols = 2, 2
    shape = slide.shapes.add_table(rows, cols, 0, 0, 1000000, 500000)
    shape.table.cell(0, 0).text = "Q1"
    shape.table.cell(0, 1).text = "100"
    presentation.core_properties.title = "Board Briefing"
    presentation.save(str(ppt_path))
    powerpoint = PowerPointParser().parse(ppt_path, metadata(ppt_path, "application/pptx"))
    assert powerpoint.metadata["slide_count"] == 1
    assert powerpoint.blocks[0].block_type is BlockType.SLIDE
    assert all(block.parent_ref == powerpoint.blocks[0].ref for block in powerpoint.blocks[1:])
    assert any(block.block_type is BlockType.TABLE for block in powerpoint.blocks)
    assert powerpoint.blocks[1].source_locators[0].shape_index is not None

    excel_path = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "指标"
    sheet.append(["名称", "Value"])
    sheet.append(["收入", 100])
    sheet["C2"] = "=B2*2"
    sheet.append(["隐藏行", 999, None, "secret-row"])
    sheet.row_dimensions[3].hidden = True
    sheet["D1"] = "隐藏列"
    sheet["D2"] = "secret-column"
    sheet.column_dimensions["D"].hidden = True
    hidden = workbook.create_sheet("隐藏")
    hidden.sheet_state = "hidden"
    workbook.save(excel_path)
    excel = ExcelParser().parse(excel_path, metadata(excel_path, "application/xlsx"))
    assert excel.metadata["visible_sheet_count"] == 1
    logical_table = next(
        block for block in excel.blocks if block.block_type is BlockType.LOGICAL_TABLE
    )
    assert logical_table.source_locators[0].sheet_name == "指标"
    assert logical_table.source_locators[0].cell_range == "A1:C2"
    assert logical_table.attributes["formulas"][1][2] == "=B2*2"
    assert "secret" not in logical_table.original_text


def test_mineru_content_list_mapping_and_zip_validation(tmp_path: Path) -> None:
    pdf_path = tmp_path / "report.pdf"
    pdf_path.write_bytes(b"%PDF-1.7\nfixture")
    file_metadata = metadata(pdf_path, "application/pdf")
    content_list = [
        {"type": "header", "text": "discard", "page_idx": 0, "bbox": [0, 0, 1, 1]},
        {"type": "text", "text": "年度报告", "text_level": 1, "page_idx": 0, "bbox": [1, 2, 3, 4]},
        {"type": "text", "text": "Revenue 收入", "page_idx": 0, "bbox": [1, 5, 3, 8]},
        {
            "type": "table",
            "table_body": "<table><tr><td>A</td><td>B</td></tr></table>",
            "table_caption": ["表一"],
            "page_idx": 1,
            "bbox": [10, 20, 30, 40],
        },
        {"type": "equation", "text": "x=1", "text_format": "latex", "page_idx": 1},
        {"type": "code", "sub_type": "algorithm", "code_body": "return 1", "page_idx": 1},
        {"type": "list", "list_items": ["a", "b"], "page_idx": 1},
        {"type": "image", "image_caption": ["图一"], "page_idx": 1},
        {"type": "page_footnote", "text": "脚注", "page_idx": 1},
    ]
    artifact = MinerUParser.from_content_list(content_list, file_metadata)
    assert artifact.title == "年度报告"
    assert artifact.metadata["page_count"] == 2
    types = [block.block_type for block in artifact.blocks]
    assert BlockType.PAGE in types
    assert BlockType.TABLE in types
    assert BlockType.FORMULA in types
    assert BlockType.CODE in types
    assert BlockType.LIST in types
    assert BlockType.CAPTION in types
    assert BlockType.FOOTNOTE in types
    assert "discard" not in " ".join(block.original_text for block in artifact.blocks)

    payload = io.BytesIO()
    with zipfile.ZipFile(payload, "w") as archive:
        archive.writestr("report/report_content_list.json", json.dumps(content_list))
    extracted, names = MinerUParser._extract_content_list(payload.getvalue())
    assert extracted == content_list
    assert names == ["report/report_content_list.json"]
    with pytest.raises(ParseError, match="valid ZIP"):
        MinerUParser._extract_content_list(b"not-a-zip")
    missing_payload = io.BytesIO()
    with zipfile.ZipFile(missing_payload, "w") as archive:
        archive.writestr("report.md", "text")
    with pytest.raises(ParseError, match="content_list"):
        MinerUParser._extract_content_list(missing_payload.getvalue())

    ppt_metadata = FileMetadata(
        filename="briefing.pptx",
        mime_type="application/pptx",
        file_size=100,
        sha256="b" * 64,
    )
    ppt = MinerUParser.from_content_list(content_list[1:3], ppt_metadata)
    assert ppt.metadata["slide_count"] == 1
    assert ppt.blocks[0].block_type is BlockType.SLIDE
    assert ppt.blocks[1].source_locators[0].slide_number == 1


def test_mineru_precision_api_signed_upload_poll_and_download(tmp_path: Path) -> None:
    pdf_path = tmp_path / "report.pdf"
    pdf_bytes = b"%PDF-1.7\nprecision"
    pdf_path.write_bytes(pdf_bytes)
    file_metadata = metadata(pdf_path, "application/pdf")
    content_list = [
        {"type": "text", "text": "精准报告", "text_level": 1, "page_idx": 0},
        {"type": "text", "text": "Revenue 收入", "page_idx": 0},
    ]
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w") as output:
        output.writestr("report/report_content_list.json", json.dumps(content_list))
        output.writestr("report/full.md", "# 精准报告")

    poll_count = 0

    def handler(request: httpx.Request) -> httpx.Response:
        nonlocal poll_count
        if request.url.path == "/api/v4/file-urls/batch":
            assert request.method == "POST"
            assert request.headers["Authorization"] == "Bearer test-token"
            body = json.loads(request.content)
            assert body["model_version"] == "vlm"
            assert body["files"] == [{"name": "report.pdf", "data_id": "a" * 64, "is_ocr": True}]
            return httpx.Response(
                200,
                json={
                    "code": 0,
                    "data": {
                        "batch_id": "batch-1",
                        "file_urls": ["https://upload.example/report.pdf"],
                    },
                    "trace_id": "trace-submit",
                },
            )
        if request.url.host == "upload.example":
            assert request.method == "PUT"
            assert "Authorization" not in request.headers
            assert request.headers["Content-Length"] == str(len(pdf_bytes))
            assert request.read() == pdf_bytes
            return httpx.Response(200)
        if request.url.path == "/api/v4/extract-results/batch/batch-1":
            assert request.headers["Authorization"] == "Bearer test-token"
            poll_count += 1
            state = "running" if poll_count == 1 else "done"
            result: dict[str, object] = {
                "file_name": "report.pdf",
                "data_id": "a" * 64,
                "state": state,
                "err_msg": "",
            }
            if state == "done":
                result["full_zip_url"] = "https://cdn.example/result.zip"
            return httpx.Response(
                200,
                json={
                    "code": 0,
                    "data": {"batch_id": "batch-1", "extract_result": [result]},
                    "trace_id": "trace-result",
                },
            )
        if request.url.host == "cdn.example":
            assert "Authorization" not in request.headers
            return httpx.Response(200, content=archive.getvalue())
        raise AssertionError(f"Unexpected request: {request.method} {request.url}")

    parser = MinerUParser(
        base_url="https://mineru.net/api/v4",
        token="test-token",
        timeout_seconds=5,
        poll_interval_seconds=0.001,
        model_version="vlm",
        transport=httpx.MockTransport(handler),
    )
    artifact = parser.parse(pdf_path, file_metadata)

    assert "test-token" not in json.dumps(parser.config_snapshot)
    assert artifact.parser_name == "mineru-precision"
    assert artifact.title == "精准报告"
    assert artifact.metadata["mineru_batch_id"] == "batch-1"
    assert artifact.metadata["mineru_trace_id"] == "trace-result"
    assert artifact.metadata["mineru_model_version"] == "vlm"
    assert poll_count == 2


def test_mineru_precision_html_and_explicit_auth_failures(tmp_path: Path) -> None:
    html_path = tmp_path / "policy.htm"
    html_path.write_text("<p>original</p>", encoding="utf-8")
    file_metadata = metadata(html_path, "text/html")
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w") as output:
        output.writestr(
            "policy/main.html",
            "<html><head><title>制度</title></head><body><h1>Policy</h1><p>正文</p></body></html>",
        )

    def handler(request: httpx.Request) -> httpx.Response:
        if request.url.path == "/api/v4/file-urls/batch":
            body = json.loads(request.content)
            assert body["model_version"] == "MinerU-HTML"
            assert body["files"][0]["name"] == "policy.html"
            return httpx.Response(
                200,
                json={
                    "code": 0,
                    "data": {
                        "batch_id": "html-batch",
                        "file_urls": ["https://upload.example/policy.html"],
                    },
                },
            )
        if request.url.host == "upload.example":
            return httpx.Response(200)
        if request.url.path.endswith("/html-batch"):
            return httpx.Response(
                200,
                json={
                    "code": 0,
                    "data": {
                        "extract_result": [
                            {
                                "data_id": "a" * 64,
                                "state": "done",
                                "full_zip_url": "https://cdn.example/html.zip",
                            }
                        ]
                    },
                },
            )
        if request.url.host == "cdn.example":
            return httpx.Response(200, content=archive.getvalue())
        raise AssertionError(f"Unexpected request: {request.method} {request.url}")

    parser = MinerUParser(
        base_url="https://mineru.net/api/v4",
        token="test-token",
        timeout_seconds=5,
        poll_interval_seconds=0.001,
        model_version="vlm",
        transport=httpx.MockTransport(handler),
    )
    artifact = parser.parse(html_path, file_metadata)
    assert artifact.title == "制度"
    assert artifact.metadata["mineru_model_version"] == "MinerU-HTML"
    assert any(block.original_text == "正文" for block in artifact.blocks)

    missing_token = MinerUParser(
        base_url="https://mineru.net/api/v4",
        token=None,
        timeout_seconds=5,
        poll_interval_seconds=1,
        model_version="vlm",
    )
    with pytest.raises(ParseError) as missing_error:
        missing_token.parse(html_path, file_metadata)
    assert missing_error.value.code == "MINERU_TOKEN_MISSING"

    def auth_handler(_request: httpx.Request) -> httpx.Response:
        return httpx.Response(200, json={"code": "A0202", "msg": "Token invalid"})

    invalid_token = MinerUParser(
        base_url="https://mineru.net/api/v4",
        token="bad-token",
        timeout_seconds=5,
        poll_interval_seconds=1,
        model_version="vlm",
        transport=httpx.MockTransport(auth_handler),
    )
    with pytest.raises(ParseError) as auth_error:
        invalid_token.parse(html_path, file_metadata)
    assert auth_error.value.code == "MINERU_AUTH_FAILED"
    assert auth_error.value.retryable is False


def test_router_and_canonicalizer_are_deterministic(tmp_path: Path) -> None:
    source = tmp_path / "source.txt"
    source.write_text("hello", encoding="utf-8")
    file_metadata = metadata(source, "text/plain")
    parser = PlainTextParser()
    router = ParserRouter([parser])
    assert router.select(source, file_metadata) is parser

    bad_pdf = tmp_path / "bad.pdf"
    bad_pdf.write_bytes(b"not-pdf")
    with pytest.raises(ParseError, match="signature"):
        router.select(bad_pdf, metadata(bad_pdf, "application/pdf"))
    with pytest.raises(ParseError, match="No parser"):
        ParserRouter([]).select(source, file_metadata)

    artifact = ParseArtifact(
        parser_name="fixture",
        parser_version="1",
        parser_mode="test",
        title="Title",
        blocks=[
            ParsedBlock(
                ref="h1",
                block_type=BlockType.HEADING,
                original_text=" 标题 ",
                attributes={"level": 1},
            ),
            ParsedBlock(
                ref="p1", block_type=BlockType.PARAGRAPH, original_text="中文 English\u00a0 text"
            ),
        ],
    )
    document_id = uuid.uuid4()
    version_id = uuid.uuid4()
    first = Canonicalizer().convert(
        artifact=artifact, document_id=document_id, version_id=version_id
    )
    second = Canonicalizer().convert(
        artifact=artifact, document_id=document_id, version_id=version_id
    )
    assert first == second
    assert first.language == "zh-en-mixed"
    assert first.blocks[2].heading_path == ["标题"]
    assert first.blocks[2].original_text != first.blocks[2].normalized_text
    assert first.blocks[2].token_count > 0

    with pytest.raises(ValidationError, match="unique"):
        ParseArtifact(
            parser_name="bad",
            parser_version="1",
            parser_mode="test",
            blocks=[
                ParsedBlock(ref="same", block_type=BlockType.PARAGRAPH),
                ParsedBlock(ref="same", block_type=BlockType.PARAGRAPH),
            ],
        )


def test_legacy_office_conversion_has_success_and_explicit_failures(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "legacy.doc"
    source.write_bytes(bytes.fromhex("D0CF11E0A1B11AE1") + b"fixture")
    file_metadata = metadata(source, "application/msword")
    parser = LegacyOfficeParser("soffice", {".docx": WordParser()})

    monkeypatch.setattr(shutil, "which", lambda _value: "/fake/soffice")

    def successful_conversion(
        arguments: list[str], **_kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        output_directory = Path(arguments[arguments.index("--outdir") + 1])
        converted = WordDocument()
        converted.add_heading("Converted", level=1)
        converted.save(str(output_directory / "legacy.docx"))
        return subprocess.CompletedProcess(arguments, 0, "converted", "")

    monkeypatch.setattr(subprocess, "run", successful_conversion)
    artifact = parser.parse(source, file_metadata)
    assert artifact.parser_name == "libreoffice-converter+python-docx"
    assert artifact.metadata["converted_from"] == ".doc"

    monkeypatch.setattr(shutil, "which", lambda _value: None)
    with pytest.raises(ParseError, match="LibreOffice is required"):
        parser.parse(source, file_metadata)

    monkeypatch.setattr(shutil, "which", lambda _value: "/fake/soffice")
    monkeypatch.setattr(
        subprocess,
        "run",
        lambda arguments, **_kwargs: subprocess.CompletedProcess(
            arguments, 1, "", "conversion failed"
        ),
    )
    with pytest.raises(ParseError, match="conversion failed"):
        parser.parse(source, file_metadata)
