"""Lightweight deterministic parsers for text, web and OOXML documents."""

import re
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, ClassVar

from bs4 import BeautifulSoup, Tag
from docx import Document as WordDocument
from docx.table import Table as WordTable
from docx.text.paragraph import Paragraph as WordParagraph
from markdown_it import MarkdownIt
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation

from robust_rag.parsing.base import FileMetadata, ParseError, Parser
from robust_rag.parsing.schemas import (
    BlockType,
    ParseArtifact,
    ParsedBlock,
    SourceLocator,
    SourceType,
)


def _ref(prefix: str, index: int) -> str:
    return f"{prefix}-{index:05d}"


class PlainTextParser:
    name = "plain-text"
    version = "1.0.0"
    mode = "native"

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension == ".txt" and metadata.mime_type == "text/plain"

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        text = source_path.read_text(encoding="utf-8-sig")
        blocks: list[ParsedBlock] = []
        for index, match in enumerate(re.finditer(r"\S(?:.*?\S)?(?=\n\s*\n|\Z)", text, re.S), 1):
            value = match.group(0)
            line_start = text.count("\n", 0, match.start()) + 1
            line_end = line_start + value.count("\n")
            blocks.append(
                ParsedBlock(
                    ref=_ref("paragraph", index),
                    block_type=BlockType.PARAGRAPH,
                    original_text=value,
                    source_locators=[
                        SourceLocator(
                            source_type=SourceType.TEXT,
                            line_start=line_start,
                            line_end=line_end,
                            char_start=match.start(),
                            char_end=match.end(),
                        )
                    ],
                )
            )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=source_path.stem,
            metadata={"filename": metadata.filename},
            blocks=blocks,
        )


class MarkdownParser:
    name = "markdown-it"
    version = "1.0.0"
    mode = "native"

    def __init__(self) -> None:
        self.markdown = MarkdownIt("commonmark", {"html": False})

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension in {".md", ".markdown"}

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        text = source_path.read_text(encoding="utf-8-sig")
        tokens = self.markdown.parse(text)
        lines = text.splitlines()
        blocks: list[ParsedBlock] = []
        index = 0
        list_stack: list[str] = []
        list_item_stack: list[str] = []

        def append_block(
            *,
            block_type: BlockType,
            value: str,
            start: int,
            end: int,
            attributes: dict[str, Any] | None = None,
            parent_ref: str | None = None,
        ) -> str:
            nonlocal index
            index += 1
            block_ref = _ref("markdown", index)
            blocks.append(
                ParsedBlock(
                    ref=block_ref,
                    parent_ref=parent_ref,
                    block_type=block_type,
                    original_text=value,
                    attributes=attributes or {},
                    source_locators=[
                        SourceLocator(
                            source_type=SourceType.MARKDOWN,
                            line_start=start + 1,
                            line_end=max(start + 1, end),
                            char_start=sum(len(line) + 1 for line in lines[:start]),
                            char_end=sum(len(line) + 1 for line in lines[:end]),
                        )
                    ],
                )
            )
            return block_ref

        for position, token in enumerate(tokens):
            start, end = token.map or [0, 0]
            if token.type in {"bullet_list_open", "ordered_list_open"}:
                list_ref = append_block(
                    block_type=BlockType.LIST,
                    value="",
                    start=start,
                    end=end,
                    attributes={"ordered": token.type == "ordered_list_open"},
                    parent_ref=list_item_stack[-1] if list_item_stack else None,
                )
                list_stack.append(list_ref)
            elif token.type in {"bullet_list_close", "ordered_list_close"}:
                list_stack.pop()
            elif token.type == "list_item_open":
                inline = next(
                    (
                        candidate
                        for candidate in tokens[position + 1 :]
                        if candidate.type == "inline"
                    ),
                    None,
                )
                item_ref = append_block(
                    block_type=BlockType.LIST_ITEM,
                    value=inline.content if inline else "",
                    start=start,
                    end=end,
                    parent_ref=list_stack[-1] if list_stack else None,
                )
                list_item_stack.append(item_ref)
            elif token.type == "list_item_close":
                list_item_stack.pop()
            elif token.type == "heading_open":
                append_block(
                    block_type=BlockType.HEADING,
                    value=tokens[position + 1].content,
                    start=start,
                    end=end,
                    attributes={"level": int(token.tag.removeprefix("h"))},
                )
            elif token.type in {"fence", "code_block"}:
                append_block(
                    block_type=BlockType.CODE,
                    value=token.content,
                    start=start,
                    end=end,
                    attributes={"language": token.info.strip() or None},
                    parent_ref=list_item_stack[-1] if list_item_stack else None,
                )
            elif token.type == "paragraph_open" and not list_item_stack:
                append_block(
                    block_type=BlockType.PARAGRAPH,
                    value=tokens[position + 1].content,
                    start=start,
                    end=end,
                )
        title = next(
            (block.original_text for block in blocks if block.block_type is BlockType.HEADING),
            source_path.stem,
        )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=title,
            metadata={"filename": metadata.filename},
            blocks=blocks,
        )


class HtmlParser:
    name = "beautifulsoup"
    version = "1.0.0"
    mode = "native"
    discarded_tags: ClassVar[set[str]] = {
        "script",
        "style",
        "nav",
        "footer",
        "aside",
        "noscript",
        "template",
    }

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension in {".html", ".htm"} and metadata.mime_type == "text/html"

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        soup = BeautifulSoup(source_path.read_text(encoding="utf-8-sig"), "html.parser")
        for node in soup.find_all(self.discarded_tags):
            node.decompose()
        blocks: list[ParsedBlock] = []
        tags = soup.find_all(
            [
                "h1",
                "h2",
                "h3",
                "h4",
                "h5",
                "h6",
                "p",
                "ul",
                "ol",
                "blockquote",
                "pre",
                "table",
            ]
        )
        index = 0
        for tag in tags:
            if tag.find_parent(["table", "pre"]) and tag.name not in {"table", "pre"}:
                continue
            if tag.name not in {"ul", "ol"} and tag.find_parent(["ul", "ol"]):
                continue
            if tag.name in {"ul", "ol"} and tag.find_parent(["ul", "ol"]):
                continue
            value = tag.get_text(" ", strip=True)
            if not value:
                continue
            index += 1
            block_ref = _ref("html", index)
            if tag.name and tag.name.startswith("h"):
                block_type = BlockType.HEADING
                attributes: dict[str, Any] = {"level": int(tag.name[1])}
            else:
                block_type = {
                    "ul": BlockType.LIST,
                    "ol": BlockType.LIST,
                    "blockquote": BlockType.QUOTE,
                    "pre": BlockType.CODE,
                    "table": BlockType.TABLE,
                }.get(tag.name or "", BlockType.PARAGRAPH)
                attributes = {"ordered": tag.name == "ol"} if tag.name in {"ul", "ol"} else {}
            links = [
                {"text": link.get_text(" ", strip=True), "url": link.get("href")}
                for link in tag.find_all("a", href=True)
            ]
            if links:
                attributes["links"] = links
            blocks.append(
                ParsedBlock(
                    ref=block_ref,
                    block_type=block_type,
                    original_text="" if tag.name in {"ul", "ol"} else value,
                    attributes=attributes,
                    source_locators=[
                        SourceLocator(source_type=SourceType.HTML, dom_path=self._dom_path(tag))
                    ],
                )
            )
            if tag.name in {"ul", "ol"}:
                for item in tag.find_all("li", recursive=False):
                    item_text = " ".join(item.find_all(string=True, recursive=False)).strip()
                    if not item_text:
                        item_text = item.get_text(" ", strip=True)
                    index += 1
                    blocks.append(
                        ParsedBlock(
                            ref=_ref("html", index),
                            parent_ref=block_ref,
                            block_type=BlockType.LIST_ITEM,
                            original_text=item_text,
                            source_locators=[
                                SourceLocator(
                                    source_type=SourceType.HTML,
                                    dom_path=self._dom_path(item),
                                )
                            ],
                        )
                    )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=soup.title.string.strip()
            if soup.title and soup.title.string
            else source_path.stem,
            metadata={"filename": metadata.filename},
            blocks=blocks,
        )

    @staticmethod
    def _dom_path(tag: Tag) -> str:
        parts: list[str] = []
        current: Tag | None = tag
        while current is not None and current.name != "[document]":
            siblings = [
                s
                for s in current.previous_siblings
                if isinstance(s, Tag) and s.name == current.name
            ]
            parts.append(f"{current.name}:nth-of-type({len(siblings) + 1})")
            parent = current.parent
            current = parent if isinstance(parent, Tag) else None
        return " > ".join(reversed(parts))


class WordParser:
    name = "python-docx"
    version = "1.0.0"
    mode = "native"

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension == ".docx"

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        document = WordDocument(str(source_path))
        blocks: list[ParsedBlock] = []
        index = 0
        paragraph_index = 0
        table_index = 0
        for child in document.element.body.iterchildren():
            if child.tag.endswith("}p"):
                paragraph_index += 1
                paragraph = WordParagraph(child, document)
                text = paragraph.text.strip()
                if not text:
                    continue
                index += 1
                style = paragraph.style.name if paragraph.style else ""
                heading_match = re.match(r"Heading\s+(\d+)", style, re.I)
                is_list = "List" in style
                attributes: dict[str, Any] = {
                    "paragraph_index": paragraph_index,
                    "style": style,
                }
                if heading_match:
                    block_type = BlockType.HEADING
                    attributes["level"] = int(heading_match.group(1))
                else:
                    block_type = BlockType.LIST_ITEM if is_list else BlockType.PARAGRAPH
                blocks.append(
                    ParsedBlock(
                        ref=_ref("word", index),
                        block_type=block_type,
                        original_text=text,
                        attributes=attributes,
                        source_locators=[
                            SourceLocator(
                                source_type=SourceType.WORD,
                                paragraph_index=paragraph_index,
                            )
                        ],
                    )
                )
            elif child.tag.endswith("}tbl"):
                table_index += 1
                table = WordTable(child, document)
                index += 1
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                blocks.append(
                    ParsedBlock(
                        ref=_ref("word", index),
                        block_type=BlockType.TABLE,
                        original_text="\n".join("\t".join(row) for row in rows),
                        attributes={"table_index": table_index, "rows": rows},
                        source_locators=[
                            SourceLocator(source_type=SourceType.WORD, table_index=table_index)
                        ],
                    )
                )
        for footnote_id, text in self._footnotes(source_path):
            index += 1
            blocks.append(
                ParsedBlock(
                    ref=_ref("word", index),
                    block_type=BlockType.FOOTNOTE,
                    original_text=text,
                    attributes={"footnote_id": footnote_id},
                    source_locators=[SourceLocator(source_type=SourceType.WORD)],
                )
            )
        title = document.core_properties.title or next(
            (b.original_text for b in blocks if b.block_type is BlockType.HEADING), source_path.stem
        )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=title,
            metadata={"filename": metadata.filename},
            blocks=blocks,
        )

    @staticmethod
    def _footnotes(source_path: Path) -> list[tuple[int, str]]:
        namespace = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        try:
            with zipfile.ZipFile(source_path) as archive:
                payload = archive.read("word/footnotes.xml")
        except KeyError:
            return []
        root = ET.fromstring(payload)
        output: list[tuple[int, str]] = []
        for footnote in root.findall(f"{{{namespace}}}footnote"):
            footnote_id = int(footnote.attrib.get(f"{{{namespace}}}id", "-1"))
            if footnote_id < 1:
                continue
            text = "".join(node.text or "" for node in footnote.iter(f"{{{namespace}}}t")).strip()
            if text:
                output.append((footnote_id, text))
        return output


class PowerPointParser:
    name = "python-pptx"
    version = "1.0.0"
    mode = "native"

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension == ".pptx"

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        presentation = Presentation(str(source_path))
        blocks: list[ParsedBlock] = []
        index = 0
        for slide_number, slide in enumerate(presentation.slides, 1):
            index += 1
            slide_ref = _ref("slide", slide_number)
            blocks.append(
                ParsedBlock(
                    ref=slide_ref,
                    block_type=BlockType.SLIDE,
                    attributes={"slide_number": slide_number},
                    source_locators=[
                        SourceLocator(source_type=SourceType.POWERPOINT, slide_number=slide_number)
                    ],
                )
            )
            for shape_index, shape in enumerate(slide.shapes, 1):
                if getattr(shape, "has_table", False):
                    index += 1
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    blocks.append(
                        ParsedBlock(
                            ref=_ref("ppt", index),
                            parent_ref=slide_ref,
                            block_type=BlockType.TABLE,
                            original_text="\n".join("\t".join(row) for row in rows),
                            attributes={"shape_index": shape_index, "rows": rows},
                            source_locators=[
                                SourceLocator(
                                    source_type=SourceType.POWERPOINT,
                                    slide_number=slide_number,
                                    shape_index=shape_index,
                                )
                            ],
                        )
                    )
                elif getattr(shape, "has_text_frame", False) and shape.text.strip():
                    index += 1
                    is_title = shape == slide.shapes.title
                    blocks.append(
                        ParsedBlock(
                            ref=_ref("ppt", index),
                            parent_ref=slide_ref,
                            block_type=BlockType.HEADING if is_title else BlockType.PARAGRAPH,
                            original_text=shape.text.strip(),
                            attributes={
                                "shape_index": shape_index,
                                **({"level": 1} if is_title else {}),
                            },
                            source_locators=[
                                SourceLocator(
                                    source_type=SourceType.POWERPOINT,
                                    slide_number=slide_number,
                                    shape_index=shape_index,
                                )
                            ],
                        )
                    )
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
            if notes:
                index += 1
                blocks.append(
                    ParsedBlock(
                        ref=_ref("ppt", index),
                        parent_ref=slide_ref,
                        block_type=BlockType.NOTE,
                        original_text=notes,
                        attributes={"speaker_notes": True},
                        source_locators=[
                            SourceLocator(
                                source_type=SourceType.POWERPOINT, slide_number=slide_number
                            )
                        ],
                    )
                )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=presentation.core_properties.title or source_path.stem,
            metadata={"filename": metadata.filename, "slide_count": len(presentation.slides)},
            blocks=blocks,
        )


class ExcelParser:
    name = "openpyxl"
    version = "1.0.0"
    mode = "native"

    def can_handle(self, metadata: FileMetadata) -> bool:
        return metadata.extension == ".xlsx"

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        values_workbook = load_workbook(source_path, data_only=True, read_only=False)
        formula_workbook = load_workbook(source_path, data_only=False, read_only=False)
        blocks: list[ParsedBlock] = []
        index = 0
        included_sheets = 0
        for sheet in values_workbook.worksheets:
            if sheet.sheet_state != "visible":
                continue
            included_sheets += 1
            index += 1
            sheet_ref = _ref("sheet", included_sheets)
            blocks.append(
                ParsedBlock(
                    ref=sheet_ref,
                    block_type=BlockType.SHEET,
                    original_text=sheet.title,
                    attributes={"sheet_index": values_workbook.sheetnames.index(sheet.title) + 1},
                    source_locators=[
                        SourceLocator(source_type=SourceType.EXCEL, sheet_name=sheet.title)
                    ],
                )
            )
            rows: list[list[Any]] = []
            formula_rows: list[list[Any]] = []
            visible_columns = [
                column
                for column in range(1, sheet.max_column + 1)
                if not sheet.column_dimensions[get_column_letter(column)].hidden
            ]
            if not visible_columns:
                continue
            visible_row_numbers: list[int] = []
            for row_number in range(1, sheet.max_row + 1):
                if sheet.row_dimensions[row_number].hidden:
                    continue
                visible_row_numbers.append(row_number)
                rows.append([sheet.cell(row_number, column).value for column in visible_columns])
                formula_rows.append(
                    [
                        formula_workbook[sheet.title].cell(row_number, column).value
                        for column in visible_columns
                    ]
                )
            while rows and not any(value is not None for value in rows[-1]):
                rows.pop()
                formula_rows.pop()
            if not rows:
                continue
            first_row = visible_row_numbers[0]
            last_row = visible_row_numbers[len(rows) - 1]
            first_column = visible_columns[0]
            last_column = visible_columns[-1]
            cell_range = (
                f"{get_column_letter(first_column)}{first_row}:"
                f"{get_column_letter(last_column)}{last_row}"
            )
            index += 1
            blocks.append(
                ParsedBlock(
                    ref=_ref("excel", index),
                    parent_ref=sheet_ref,
                    block_type=BlockType.LOGICAL_TABLE,
                    original_text="\n".join(
                        "\t".join("" if value is None else str(value) for value in row)
                        for row in rows
                    ),
                    attributes={
                        "display_values": rows,
                        "formulas": formula_rows,
                        "hidden_rows_excluded": True,
                        "hidden_columns_excluded": True,
                    },
                    source_locators=[
                        SourceLocator(
                            source_type=SourceType.EXCEL,
                            sheet_name=sheet.title,
                            cell_range=cell_range,
                        )
                    ],
                )
            )
        return ParseArtifact(
            parser_name=self.name,
            parser_version=self.version,
            parser_mode=self.mode,
            title=source_path.stem,
            metadata={"filename": metadata.filename, "visible_sheet_count": included_sheets},
            blocks=blocks,
        )


class LegacyOfficeParser:
    """Convert legacy OLE Office files with LibreOffice, then use an OOXML parser."""

    name = "libreoffice-converter"
    version = "1.0.0"
    mode = "legacy-conversion"
    conversions: ClassVar[dict[str, str]] = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}

    def __init__(self, executable: str, delegates: dict[str, Parser]) -> None:
        self.executable = executable
        self.delegates = delegates

    def can_handle(self, metadata: FileMetadata) -> bool:
        target_extension = self.conversions.get(metadata.extension)
        return target_extension is not None and target_extension in self.delegates

    def parse(self, source_path: Path, metadata: FileMetadata) -> ParseArtifact:
        executable = shutil.which(self.executable)
        if executable is None:
            raise ParseError(
                "LIBREOFFICE_UNAVAILABLE", "LibreOffice is required for legacy Office files"
            )
        target_extension = self.conversions[metadata.extension]
        with tempfile.TemporaryDirectory(prefix="robust-rag-office-") as temporary_directory:
            output_directory = Path(temporary_directory)
            result = subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    target_extension.removeprefix("."),
                    "--outdir",
                    str(output_directory),
                    str(source_path),
                ],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            converted = output_directory / f"{source_path.stem}{target_extension}"
            if result.returncode != 0 or not converted.exists():
                raise ParseError(
                    "OFFICE_CONVERSION_FAILED",
                    result.stderr.strip() or "LibreOffice conversion failed",
                )
            delegate = self.delegates[target_extension]
            converted_metadata = FileMetadata(
                filename=converted.name,
                mime_type="application/zip",
                file_size=converted.stat().st_size,
                sha256=metadata.sha256,
            )
            artifact = delegate.parse(converted, converted_metadata)
            return artifact.model_copy(
                update={
                    "parser_name": f"{self.name}+{artifact.parser_name}",
                    "parser_mode": self.mode,
                    "metadata": {**artifact.metadata, "converted_from": metadata.extension},
                }
            )
